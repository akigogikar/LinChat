from __future__ import annotations

import json
from typing import List, Optional

import openai
import pandas as pd
from pydantic import BaseModel
from weasyprint import HTML
import markdown as md
from pptx import Presentation
from pptx.util import Inches
import os

from .config import _get_api_key, _get_model, _read_config


class Summary(BaseModel):
    title: str
    bullets: List[str]


class Table(BaseModel):
    columns: List[str]
    rows: List[List[str]]


class Slide(BaseModel):
    title: str
    bullets: List[str] = []
    table: Optional[Table] = None


class SlideDeck(BaseModel):
    slides: List[Slide]




# LLM generation functions

def _call_llm(prompt: str, schema: dict, fn_name: str) -> dict:
    openai.api_key = _get_api_key()
    openai.base_url = "https://openrouter.ai/api/v1"
    completion = openai.ChatCompletion.create(
        model=_get_model(),
        messages=[{"role": "user", "content": prompt}],
        functions=[{"name": fn_name, "parameters": schema}],
        function_call={"name": fn_name},
    )
    args = completion.choices[0].message["function_call"]["arguments"]
    return json.loads(args)


def generate_summary(prompt: str) -> Summary:
    data = _call_llm(prompt, Summary.schema(), "generate_summary")
    return Summary(**data)


def generate_table(prompt: str) -> Table:
    data = _call_llm(prompt, Table.schema(), "generate_table")
    return Table(**data)


def generate_slide_deck(prompt: str) -> SlideDeck:
    data = _call_llm(prompt, SlideDeck.schema(), "generate_slide_deck")
    return SlideDeck(**data)


# Export helpers

def markdown_to_pdf(text: str, output_path: str) -> None:
    html = md.markdown(text)
    HTML(string=html).write_pdf(output_path)


def html_to_pdf(html: str, output_path: str) -> None:
    HTML(string=html).write_pdf(output_path)


def dataframe_to_excel(df: pd.DataFrame, output_path: str) -> None:
    df.to_excel(output_path, index=False)


def slide_deck_to_pptx(deck: SlideDeck, output_path: str) -> None:
    prs = Presentation()
    for slide_data in deck.slides:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_data.title
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        for bullet in slide_data.bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
        if slide_data.table:
            rows = len(slide_data.table.rows)
            cols = len(slide_data.table.columns)
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(0.8 + 0.2 * rows)
            table_shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
            table = table_shape.table
            for c, name in enumerate(slide_data.table.columns):
                table.cell(0, c).text = name
            for r, row in enumerate(slide_data.table.rows, start=1):
                for c, cell in enumerate(row):
                    table.cell(r, c).text = str(cell)
    prs.save(output_path)
