from typing import List, Dict
from pptx import Presentation


def parse_pptx(path: str) -> List[Dict]:
    """Parse PowerPoint file and return text chunks per slide."""
    prs = Presentation(path)
    chunks = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        text = "\n".join(texts).strip()
        if text:
            chunks.append({"page": idx, "text": text})
    return chunks
