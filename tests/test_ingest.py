import os
import tempfile
from app.ingest import parse_file, parse_pdf, parse_docx, parse_excel, parse_pptx
from docx import Document
from openpyxl import Workbook
from pptx import Presentation


def _create_docx(path: str):
    doc = Document()
    doc.add_paragraph("Hello Docx")
    doc.save(path)


def _create_excel(path: str):
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Hello Excel"
    wb.save(path)


def _create_pptx(path: str):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello PPTX"
    prs.save(path)


def test_parse_docx(tmp_path):
    f = tmp_path / "test.docx"
    _create_docx(f)
    chunks = parse_docx(str(f))
    assert chunks and chunks[0]["text"] == "Hello Docx"


def test_parse_excel(tmp_path):
    f = tmp_path / "test.xlsx"
    _create_excel(f)
    chunks = parse_excel(str(f))
    assert chunks and "Hello Excel" in chunks[0]["text"]


def test_parse_pptx(tmp_path):
    f = tmp_path / "test.pptx"
    _create_pptx(f)
    chunks = parse_pptx(str(f))
    assert chunks and "Hello PPTX" in chunks[0]["text"]


def test_parse_pdf():
    path = os.path.join(os.path.dirname(__file__), "fixtures", "hello.pdf")
    chunks = parse_pdf(path)
    assert any("Hello" in c["text"] for c in chunks)


def test_parse_file_dispatch(tmp_path):
    f = tmp_path / "dispatch.docx"
    _create_docx(f)
    chunks = parse_file(str(f))
    assert chunks and chunks[0]["text"] == "Hello Docx"
